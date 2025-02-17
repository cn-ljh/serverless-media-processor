import os
import io
import logging
from enum import Enum
from typing import List, Optional, Tuple
import subprocess
import tempfile
from PIL import Image
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import pandas as pd
from b64encoder_decoder import custom_b64decode

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class ProcessingError(Exception):
    """Custom exception for processing errors"""
    def __init__(self, status_code: int, detail: str):
        self.status_code = status_code
        self.detail = detail
        super().__init__(detail)

class SourceFormat(str, Enum):
    # Word formats
    DOC = "doc"
    DOCX = "docx"
    WPS = "wps"
    WPSS = "wpss"
    DOCM = "docm"
    DOTM = "dotm"
    DOT = "dot"
    DOTX = "dotx"
    HTML = "html"
    
    # PPT formats
    PPTX = "pptx"
    PPT = "ppt"
    POT = "pot"
    POTX = "potx"
    PPS = "pps"
    PPSX = "ppsx"
    DPS = "dps"
    DPT = "dpt"
    PPTM = "pptm"
    POTM = "potm"
    PPSM = "ppsm"
    DPSS = "dpss"
    
    # Excel formats
    XLS = "xls"
    XLT = "xlt"
    ET = "et"
    ETT = "ett"
    XLSX = "xlsx"
    XLTX = "xltx"
    CSV = "csv"
    XLSB = "xlsb"
    XLSM = "xlsm"
    XLTM = "xltm"
    ETS = "ets"
    
    # PDF format
    PDF = "pdf"
    
    # Text format
    TXT = "txt"

class TargetFormat(str, Enum):
    PDF = "pdf"
    PNG = "png"
    JPG = "jpg"
    TXT = "txt"

def parse_pages_param(pages_str: Optional[str]) -> List[int]:
    """
    Parse pages parameter from base64 encoded string
    Example: "1,2,4-10" -> [1, 2, 4, 5, 6, 7, 8, 9, 10]
    """
    if not pages_str:
        return []
        
    decoded = custom_b64decode(pages_str)
    pages = []
    
    for part in decoded.split(','):
        if '-' in part:
            start, end = map(int, part.split('-'))
            pages.extend(range(start, end + 1))
        else:
            pages.append(int(part))
            
    return sorted(list(set(pages)))

def convert_to_pdf(input_path: str, output_path: str, source_format: SourceFormat):
    """Convert document to PDF using LibreOffice"""
    try:
        logger.info(f"Converting {source_format} to PDF: {input_path} -> {output_path}")
        
        # For PDF files, copy directly
        if source_format == SourceFormat.PDF:
            with open(input_path, 'rb') as src, open(output_path, 'wb') as dst:
                dst.write(src.read())
            return

        # Use LibreOffice for conversion
        cmd = [
            'libreoffice7.6',
            '--headless',
            '--invisible',
            '--nodefault',
            '--view',
            '--nolockcheck',
            '--nologo',
            '--norestore',
            '--convert-to', 'pdf',
            '--outdir', os.path.dirname(output_path),
            input_path
        ]
        
        logger.info(f"Executing command: {' '.join(cmd)}")
        result = subprocess.run(cmd, capture_output=True, text=True)
        
        if result.returncode != 0:
            error_msg = f"LibreOffice conversion failed: {result.stderr}"
            logger.error(error_msg)
            raise ProcessingError(status_code=500, detail=error_msg)
        
        # Rename output file
        temp_pdf = os.path.join(
            os.path.dirname(output_path),
            os.path.splitext(os.path.basename(input_path))[0] + '.pdf'
        )
        os.rename(temp_pdf, output_path)
        logger.info("PDF conversion successful")
        
    except subprocess.CalledProcessError as e:
        error_msg = f"Failed to convert to PDF: {str(e)}"
        logger.error(error_msg)
        raise ProcessingError(status_code=500, detail=error_msg)
    except Exception as e:
        error_msg = f"Unexpected error during PDF conversion: {str(e)}"
        logger.error(error_msg)
        raise ProcessingError(status_code=500, detail=error_msg)

def convert_pdf_to_images(pdf_path: str, output_path: str, pages: List[int] = None, format: str = 'PNG', dpi: int = 300):
    """Convert PDF pages to images with specified DPI"""
    try:
        logger.info(f"Converting PDF to {format} images: {pdf_path} -> {output_path}")
        pdf_document = fitz.open(pdf_path)
        images = []
        
        # If no pages specified, convert all pages
        if not pages:
            pages = list(range(1, pdf_document.page_count + 1))
        
        logger.info(f"Processing pages: {pages}")
        for page_num in pages:
            if page_num > pdf_document.page_count:
                logger.warning(f"Skipping page {page_num} as it exceeds document length")
                continue
            
            page = pdf_document[page_num - 1]
            
            # Calculate appropriate zoom factor
            zoom = dpi/72
            mat = fitz.Matrix(zoom, zoom)
            
            try:
                # Get page image
                pix = page.get_pixmap(matrix=mat, alpha=False)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                images.append(img)
                
                # Cleanup
                del pix
                
            except Exception as e:
                logger.error(f"Error processing page {page_num}: {str(e)}")
                raise ValueError(f"Failed to convert page {page_num} to image: {str(e)}")
        
        if not images:
            error_msg = "No valid pages to convert"
            logger.error(error_msg)
            raise ProcessingError(status_code=400, detail=error_msg)
        
        # Save images
        if format.upper() == 'PDF':
            logger.info("Saving images as PDF")
            images[0].save(
                output_path,
                "PDF",
                save_all=True,
                append_images=images[1:] if len(images) > 1 else [],
                resolution=dpi
            )
        else:
            if len(images) > 1:
                logger.info(f"Saving multiple pages as {format}")
                images[0].save(
                    output_path,
                    format=format.upper(),
                    save_all=True,
                    append_images=images[1:],
                    optimize=True,
                    quality=95 if format.upper() in ['JPG', 'JPEG'] else None
                )
            else:
                logger.info(f"Saving single page as {format}")
                images[0].save(
                    output_path,
                    format=format.upper(),
                    optimize=True,
                    quality=95 if format.upper() in ['JPG', 'JPEG'] else None
                )
        
        logger.info("Image conversion successful")
        
    except Exception as e:
        error_msg = f"Failed to convert PDF to images: {str(e)}"
        logger.error(error_msg)
        raise ProcessingError(status_code=500, detail=error_msg)
    finally:
        if 'pdf_document' in locals():
            pdf_document.close()

import xml.etree.ElementTree as ET

def extract_text_from_word(doc_path: str) -> str:
    """Extract text from Word document by converting to XML using LibreOffice"""
    logger.info(f"Extracting text from Word document: {doc_path}")
    
    # Create a temporary directory for intermediate files
    with tempfile.TemporaryDirectory() as temp_dir:
        try:
            # First convert to XML format
            cmd = [
                'libreoffice7.6',
                '--headless',
                '--invisible',
                '--nodefault',
                '--nolockcheck',
                '--nologo',
                '--norestore',
                '--convert-to', 'xml',  # Convert to XML format
                '--outdir', temp_dir,
                doc_path
            ]
            
            logger.info(f"Converting Word to XML: {' '.join(cmd)}")
            result = subprocess.run(cmd, capture_output=True, text=True)
            
            if result.returncode != 0:
                error_msg = f"LibreOffice conversion failed: {result.stderr}"
                logger.error(error_msg)
                raise ProcessingError(status_code=500, detail=error_msg)
            
            # Get the converted XML file path
            xml_path = os.path.join(
                temp_dir,
                os.path.splitext(os.path.basename(doc_path))[0] + '.xml'
            )
            
            if not os.path.exists(xml_path):
                error_msg = "XML file was not created by LibreOffice"
                logger.error(error_msg)
                raise ProcessingError(status_code=500, detail=error_msg)
            
            # Parse XML and extract text content with size limit
            MAX_TEXT_SIZE = 5 * 1024 * 1024  # 5MB limit
            current_size = 0
            text_parts = []
            
            # Use iterparse to process XML in chunks
            context = ET.iterparse(xml_path, events=('end',))
            
            # Track text elements in LibreOffice XML format
            for event, elem in context:
                # LibreOffice XML uses specific namespaces for text content
                if 'text' in elem.tag:
                    # Handle paragraph elements
                    if 'p' in elem.tag or 'h' in elem.tag:
                        text = ''
                        # Extract text from all nested text spans
                        for child in elem.iter():
                            if 'text' in child.tag and child.text:
                                text += child.text + ' '
                            if child.tail:
                                text += child.tail + ' '
                    
                    text = text.strip()
                    if text:
                        # Check size limit
                        text_size = len(text.encode('utf-8'))
                        if current_size + text_size > MAX_TEXT_SIZE:
                            logger.warning("Text size limit reached, truncating content")
                            break
                        
                        text_parts.append(text)
                        current_size += text_size
                
                # Clear element to free memory
                elem.clear()
            
            # Clear XML tree
            del context
            
            if not text_parts:
                logger.warning("No text content found in the document")
                return "No text content found in the document"
            
            # Join with proper spacing and structure
            return '\n\n'.join(text_parts)
            
        except ET.ParseError as e:
            error_msg = f"Failed to parse XML: {str(e)}"
            logger.error(error_msg)
            raise ProcessingError(status_code=500, detail=error_msg)
        except Exception as e:
            error_msg = f"Failed to extract text: {str(e)}"
            logger.error(error_msg)
            raise ProcessingError(status_code=500, detail=error_msg)

def extract_text_from_ppt(ppt_path: str) -> str:
    """Extract text from PowerPoint document"""
    logger.info(f"Extracting text from PowerPoint document: {ppt_path}")
    
    # Check if it's a .ppt file
    if ppt_path.lower().endswith('.ppt'):
        logger.info("Detected .ppt format, using LibreOffice for text extraction")
        # Create a temporary directory for intermediate files
        with tempfile.TemporaryDirectory() as temp_dir:
            try:
                # First convert to XML format
                cmd = [
                    'libreoffice7.6',
                    '--headless',
                    '--invisible',
                    '--nodefault',
                    '--nolockcheck',
                    '--nologo',
                    '--norestore',
                    '--convert-to', 'xml',  # Convert to XML format
                    '--outdir', temp_dir,
                    ppt_path
                ]
                
                logger.info(f"Converting PPT to XML: {' '.join(cmd)}")
                result = subprocess.run(cmd, capture_output=True, text=True)
                
                if result.returncode != 0:
                    error_msg = f"LibreOffice conversion failed: {result.stderr}"
                    logger.error(error_msg)
                    raise ProcessingError(status_code=500, detail=error_msg)
                
                # Get the converted XML file path
                xml_path = os.path.join(
                    temp_dir,
                    os.path.splitext(os.path.basename(ppt_path))[0] + '.xml'
                )
                
                if not os.path.exists(xml_path):
                    error_msg = "XML file was not created by LibreOffice"
                    logger.error(error_msg)
                    raise ProcessingError(status_code=500, detail=error_msg)
                
                # Parse XML and extract text content
                text_parts = []
                current_slide = 0
                
                # Use iterparse to process XML in chunks
                context = ET.iterparse(xml_path, events=('end',))
                
                for event, elem in context:
                    # Look for page/slide elements
                    if 'page' in elem.tag:
                        current_slide += 1
                        slide_text = [f"\n[Slide {current_slide}]"]
                        
                        # Extract text from all text elements in the slide
                        for text_elem in elem.iter():
                            if 'text' in text_elem.tag and text_elem.text:
                                text = text_elem.text.strip()
                                if text:
                                    slide_text.append(text)
                        
                        if len(slide_text) > 1:  # Only add slides with content
                            text_parts.extend(slide_text)
                        
                        # Clear element to free memory
                        elem.clear()
                
                # Clear XML tree
                del context
                
                if not text_parts:
                    logger.warning("No text content found in the presentation")
                    return "No text content found in the presentation"
                
                return '\n'.join(text_parts)
                
            except ET.ParseError as e:
                error_msg = f"Failed to parse XML: {str(e)}"
                logger.error(error_msg)
                raise ProcessingError(status_code=500, detail=error_msg)
            except Exception as e:
                error_msg = f"Failed to extract text: {str(e)}"
                logger.error(error_msg)
                raise ProcessingError(status_code=500, detail=error_msg)
    else:
        # For .pptx files, use python-pptx
        logger.info("Using python-pptx for text extraction")
        try:
            prs = Presentation(ppt_path)
            text_parts = []
            
            for i, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_text.append(f"\n[Slide {i}]")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if len(slide_text) > 1:  # Only add slides with content
                    text_parts.extend(slide_text)
            
            if not text_parts:
                logger.warning("No text content found in the presentation")
                return "No text content found in the presentation"
                
            return '\n'.join(text_parts)
            
        except Exception as e:
            error_msg = f"Failed to extract text: {str(e)}"
            logger.error(error_msg)
            raise ProcessingError(status_code=500, detail=error_msg)


def extract_text_from_excel(excel_path: str) -> str:
    """Extract text from Excel document"""
    logger.info(f"Extracting text from Excel document: {excel_path}")
    try:
        # Try reading with pandas
        df = pd.read_excel(excel_path, sheet_name=None)  # Read all sheets
        text_parts = []
        
        for sheet_name, sheet_df in df.items():
            text_parts.append(f"\n[Sheet: {sheet_name}]")
            
            # Convert sheet to string with proper formatting
            sheet_text = sheet_df.to_string(index=False)
            if sheet_text.strip():
                text_parts.append(sheet_text)
        
        return '\n\n'.join(text_parts)
        
    except Exception as e:
        logger.error(f"Error reading Excel file with pandas: {str(e)}")
        raise ProcessingError(
            status_code=500,
            detail=f"Failed to extract text from Excel file: {str(e)}"
        )

def extract_text_from_csv(csv_path: str) -> str:
    """Extract text from CSV document"""
    logger.info(f"Extracting text from CSV document: {csv_path}")
    try:
        # Try reading with pandas
        df = pd.read_csv(csv_path)
        
        # Convert to string with proper formatting
        text = df.to_string(index=False)
        return text if text.strip() else "No content found in CSV file"
        
    except Exception as e:
        logger.error(f"Error reading CSV file with pandas: {str(e)}")
        raise ProcessingError(
            status_code=500,
            detail=f"Failed to extract text from CSV file: {str(e)}"
        )


def convert_pdf_to_text(pdf_path: str, output_path: str, pages: List[int] = None):
    """Extract text from PDF"""
    try:
        logger.info(f"Converting PDF to text: {pdf_path} -> {output_path}")
        pdf_document = fitz.open(pdf_path)
        text = []
        
        # If no pages specified, convert all pages
        if not pages:
            pages = list(range(1, pdf_document.page_count + 1))
        
        logger.info(f"Processing pages: {pages}")
        for page_num in pages:
            if page_num > pdf_document.page_count:
                logger.warning(f"Skipping page {page_num} as it exceeds document length")
                continue
            
            page = pdf_document[page_num - 1]
            page_text = page.get_text()
            if page_text.strip():  # Only add non-empty pages
                text.append(f"[Page {page_num}]\n{page_text}")
        
        if text:
            logger.info("Writing extracted text to file")
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write('\n\n'.join(text))
        else:
            error_msg = "No valid pages to extract text from"
            logger.error(error_msg)
            raise ProcessingError(status_code=400, detail=error_msg)
            
        logger.info("Text extraction successful")
        
    except Exception as e:
        error_msg = f"Failed to convert PDF to text: {str(e)}"
        logger.error(error_msg)
        raise ProcessingError(status_code=500, detail=error_msg)
    finally:
        if 'pdf_document' in locals():
            pdf_document.close()

def convert_document(input_path: str, output_path: str, source_format: SourceFormat, 
                    target_format: TargetFormat, pages: List[int] = None, dpi: int = 300):
    """
    Convert document from source format to target format
    
    Args:
        input_path: Path to source document
        output_path: Path for converted document
        source_format: Source document format
        target_format: Target document format
        pages: Optional list of page numbers to convert
        dpi: Resolution for image conversion (default: 300)
        
    Supported conversions:
        - All formats -> PDF
        - All formats -> PNG/JPG (by converting to PDF first then to image)
        - Word/PPT/PDF -> TXT
    """
    logger.info(f"Converting document: {input_path} ({source_format}) -> {output_path} ({target_format})")
    
    # Create temporary directory for intermediate files
    with tempfile.TemporaryDirectory() as temp_dir:
        try:
            # Handle text extraction
            if target_format == TargetFormat.TXT:
                # Extract text based on source format
                if source_format in [SourceFormat.DOC, SourceFormat.DOCX, SourceFormat.WPS]:
                    logger.info("Extracting text directly from Word document")
                    text = extract_text_from_word(input_path)
                elif source_format in [SourceFormat.PPT, SourceFormat.PPTX, SourceFormat.DPS]:
                    logger.info("Extracting text directly from PowerPoint document")
                    text = extract_text_from_ppt(input_path)
                else:
                    # Convert to PDF first then extract text
                    if source_format != SourceFormat.PDF:
                        logger.info(f"Converting {source_format} to PDF")
                        temp_pdf = os.path.join(temp_dir, 'temp.pdf')
                        convert_to_pdf(input_path, temp_pdf, source_format)
                        pdf_path = temp_pdf
                    else:
                        pdf_path = input_path
                    
                    logger.info("Extracting text from PDF")
                    convert_pdf_to_text(pdf_path, output_path, pages)
                    return

                # Write extracted text to file
                logger.info("Writing extracted text to file")
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(text)
                return
            
            # Handle document conversion
            if target_format == TargetFormat.PDF:
                if source_format == SourceFormat.PDF and not pages:
                    logger.info("Copying PDF file directly")
                    with open(input_path, 'rb') as src, open(output_path, 'wb') as dst:
                        dst.write(src.read())
                else:
                    # Convert to PDF first if needed
                    if source_format != SourceFormat.PDF:
                        logger.info(f"Converting {source_format} to PDF")
                        temp_pdf = os.path.join(temp_dir, 'temp.pdf')
                        convert_to_pdf(input_path, temp_pdf, source_format)
                        pdf_path = temp_pdf
                    else:
                        pdf_path = input_path

                    if pages:
                        logger.info("Converting specified PDF pages")
                        convert_pdf_to_images(pdf_path, output_path, pages, 'PDF', dpi)
                    else:
                        logger.info("Copying complete PDF file")
                        with open(pdf_path, 'rb') as src, open(output_path, 'wb') as dst:
                            dst.write(src.read())
            
            elif target_format in [TargetFormat.PNG, TargetFormat.JPG]:
                # Convert to PDF first
                if source_format != SourceFormat.PDF:
                    logger.info(f"Converting {source_format} to PDF")
                    temp_pdf = os.path.join(temp_dir, 'temp.pdf')
                    convert_to_pdf(input_path, temp_pdf, source_format)
                    pdf_path = temp_pdf
                else:
                    pdf_path = input_path
                
                logger.info(f"Converting PDF to {target_format} image")
                convert_pdf_to_images(pdf_path, output_path, pages, target_format, dpi)
            
            else:
                error_msg = f"Unsupported target format: {target_format}"
                logger.error(error_msg)
                raise ProcessingError(status_code=400, detail=error_msg)
                
            logger.info("Document conversion completed successfully")
            
        except Exception as e:
            error_msg = f"Document conversion failed: {str(e)}"
            logger.error(error_msg)
            raise ProcessingError(status_code=500, detail=error_msg)
